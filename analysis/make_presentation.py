"""Generate a PowerPoint presentation summarizing 3 experiments: catdog, horse, room."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PROJ = "/home/jovyan/shares/SR006.nfs2/svgrozny/project/clear_project"
NFS3 = "/home/jovyan/shares/SR006.nfs3/svgrozny"

# Slide dimensions: widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_SLIDE = RGBColor(0x16, 0x21, 0x3e)
ACCENT = RGBColor(0x00, 0xd2, 0xff)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, bold=False,
             color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16, color=WHITE):
    """lines: list of (text, bold, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = clr or color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return txBox


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.isfile(path):
        kwargs = {}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(path, left, top, **kwargs)
    return None


def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Binary Path Space Exploration",
             font_size=44, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(3.0), Inches(11), Inches(1),
             "Exhaustive Enumeration of 2^20 Diffusion Paths\nfor Image Editing via Null-Text Inversion",
             font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
             "3 Experiments  |  1M+ images each  |  8x A100 GPUs",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def _add_rounded_box(slide, left, top, width, height, fill_rgb, text, font_size=14,
                      font_color=WHITE, bold=False, subtitle=None, sub_size=11):
    """Add a rounded rectangle with centered text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.1

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color
    p.font.name = "Calibri"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(sub_size)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER

    # Vertical centering
    tf.paragraphs[0].space_before = Pt(0)
    from pptx.oxml.ns import qn
    txBody = shape._element.txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')

    return shape


def _add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    """Add a connector arrow between two points."""
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    # Add arrowhead
    ln = connector._element.find(qn('a:ln') if connector._element.find(qn('a:ln')) is not None
                                  else './/{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    if ln is None:
        from lxml import etree
        spPr = connector._element.find(qn('p:spPr'))
        if spPr is None:
            spPr = connector._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
        ln = spPr.find(qn('a:ln'))
    if ln is not None:
        from lxml import etree
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')


def make_method_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
             "Method Overview", font_size=36, bold=True, color=ACCENT)

    # --- Top row: main pipeline flow (horizontal) ---
    BOX_W = Inches(2.6)
    BOX_H = Inches(1.4)
    ROW_Y = Inches(1.3)
    GAP = Inches(0.5)
    START_X = Inches(0.3)

    box_colors = [
        RGBColor(0x0d, 0x47, 0xa1),  # blue
        RGBColor(0x1b, 0x5e, 0x20),  # green
        RGBColor(0x4a, 0x14, 0x8c),  # purple
        RGBColor(0x00, 0x69, 0x5c),  # teal
    ]

    steps = [
        ("1. Null-Text\nInversion", "Source image -> latent space\n40 diffusion steps"),
        ("2. Binary Path\nSampling", "2^20 unique masks\neach bit x2 = 40 steps"),
        ("3. Segmentation\n& Metrics", "CLIPSeg mask + bg_clip,\nbg_ssim, fg_clip_score"),
        ("4. DINOv2\nFeatures", "384-dim embeddings\nfor all 1M images"),
    ]

    box_positions = []
    for i, (title, sub) in enumerate(steps):
        x = START_X + i * (BOX_W + GAP)
        _add_rounded_box(slide, x, ROW_Y, BOX_W, BOX_H, box_colors[i],
                          title, font_size=15, bold=True, subtitle=sub, sub_size=11)
        box_positions.append((x, ROW_Y))

    # Arrows between boxes
    for i in range(len(steps) - 1):
        x1 = box_positions[i][0] + BOX_W
        x2 = box_positions[i + 1][0]
        y_mid = ROW_Y + BOX_H // 2
        _add_arrow(slide, x1, y_mid, x2, y_mid)

    # --- Step 5 box (analysis) centered below ---
    analysis_w = Inches(3.0)
    analysis_x = START_X + int(1.5 * (BOX_W + GAP)) + (BOX_W - analysis_w) // 2
    analysis_y = ROW_Y + BOX_H + Inches(0.6)
    _add_rounded_box(slide, analysis_x, analysis_y, analysis_w, Inches(1.1),
                      RGBColor(0x88, 0x00, 0x0e), "5. Analysis",
                      font_size=15, bold=True,
                      subtitle="UMAP + HDBSCAN clustering", sub_size=11)
    # Arrow down from step 4
    _add_arrow(slide, box_positions[3][0] + BOX_W // 2, ROW_Y + BOX_H,
               analysis_x + analysis_w // 2, analysis_y)

    # --- Key Idea box at bottom ---
    idea_y = Inches(4.2)
    idea_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), idea_y, Inches(11.7), Inches(1.4)
    )
    idea_shape.fill.solid()
    idea_shape.fill.fore_color.rgb = RGBColor(0x33, 0x2b, 0x00)
    idea_shape.line.color.rgb = ORANGE
    idea_shape.line.width = Pt(2)
    idea_shape.adjustments[0] = 0.05

    tf = idea_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Idea"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = ("20-bit binary mask controls 40 diffusion steps (each bit repeated x2).\n"
               "Each bit selects source or target prompt embedding for a pair of consecutive steps.\n"
               "Full enumeration: 2^20 = 1,048,576 unique paths spanning the entire interpolation space.")
    p2.font.size = Pt(15)
    p2.font.color.rgb = WHITE
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    # --- Binary mask illustration ---
    mask_y = Inches(5.7)
    add_text(slide, Inches(0.8), mask_y, Inches(11.5), Inches(0.35),
             "20-bit mask:     0  1  1  0  0  1  0  1  1  1  0  0  1  1  0  1  0  0  1  0",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), mask_y + Inches(0.35), Inches(11.5), Inches(0.35),
             "40 steps:    0 0  1 1  1 1  0 0  0 0  1 1  0 0  1 1  1 1  1 1  0 0  0 0  1 1  1 1  0 0  1 1  0 0  0 0  1 1  0 0",
             font_size=12, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

    # Bit legend
    add_text(slide, Inches(1.5), mask_y + Inches(0.75), Inches(4), Inches(0.35),
             "0 = use source prompt embedding",
             font_size=13, color=RGBColor(0x64, 0xb5, 0xf6), alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(7), mask_y + Inches(0.75), Inches(4), Inches(0.35),
             "1 = use target prompt embedding",
             font_size=13, color=RGBColor(0xff, 0xb7, 0x4d), alignment=PP_ALIGN.CENTER)


def make_experiment_slide(prs, title, source_prompt, target_prompt, seg_prompt,
                          source_img, path_source, path_target, seg_vis,
                          metrics_text, num_images):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             title, font_size=32, bold=True, color=ACCENT)

    # Prompts
    lines = [
        (f"Source: \"{source_prompt}\"", False, LIGHT_GRAY),
        (f"Target: \"{target_prompt}\"", False, LIGHT_GRAY),
        (f"Segmentation: \"{seg_prompt}\"    |    Images: {num_images}", False, LIGHT_GRAY),
    ]
    add_multiline(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(1.2), lines, font_size=14)

    # Row 1: Source image, path_source (all-0), path_target (all-1)
    img_top = Inches(2.1)
    img_h = Inches(2.6)
    labels_y = img_top - Inches(0.3)

    add_text(slide, Inches(0.5), labels_y, Inches(3), Inches(0.3),
             "Original Image", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, source_img, Inches(0.7), img_top, height=img_h)

    add_text(slide, Inches(4.0), labels_y, Inches(3), Inches(0.3),
             "Path 0...0 (source reconstruction)", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, path_source, Inches(4.2), img_top, height=img_h)

    add_text(slide, Inches(7.5), labels_y, Inches(3), Inches(0.3),
             "Path 1...1 (full target)", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, path_target, Inches(7.7), img_top, height=img_h)

    # Segmentation
    add_text(slide, Inches(10.8), labels_y, Inches(2.5), Inches(0.3),
             "Segmentation Mask", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if seg_vis and os.path.isfile(seg_vis):
        add_image_safe(slide, seg_vis, Inches(10.5), img_top, width=Inches(2.8))

    # Metrics
    add_multiline(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(2),
                  metrics_text, font_size=16)


def make_clusters_slide(prs, title, map_clusters, map_metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             f"{title} -- UMAP & Metrics Maps", font_size=32, bold=True, color=ACCENT)

    add_image_safe(slide, map_clusters, Inches(0.3), Inches(1.0), height=Inches(6.0))
    add_image_safe(slide, map_metrics, Inches(6.3), Inches(1.0), height=Inches(6.0))


def make_grid_slide(prs, title, grid_path, subtitle="Top 20 images by combined score"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             f"{title} -- {subtitle}", font_size=28, bold=True, color=ACCENT)
    add_image_safe(slide, grid_path, Inches(0.5), Inches(0.9), height=Inches(6.3))


def make_comparison_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
             "Metrics Comparison Across Experiments", font_size=32, bold=True, color=ACCENT)

    header = ("Metric", "Cat -> Dog", "Horse -> Robot", "Table -> Aquarium")
    rows = [
        ("bg_clip_similarity", "0.848 +/- 0.059", "0.862 +/- 0.047", "0.924 +/- 0.027"),
        ("bg_ssim", "0.545 +/- 0.024", "0.623 +/- 0.042", "0.572 +/- 0.073"),
        ("fg_clip_score", "0.045 +/- 0.067", "0.120 +/- 0.051", "-0.156 +/- 0.019"),
        ("Images generated", "1,048,576", "1,048,576", "1,048,576"),
        ("HDBSCAN clusters", "52", "53", "65"),
    ]

    table_shape = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.8), Inches(1.2),
                                          Inches(11.5), Inches(4.5))
    table = table_shape.table

    # Style header
    for j, txt in enumerate(header):
        cell = table.cell(0, j)
        cell.text = txt
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x0d, 0x47, 0xa1)

    # Style rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(15)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
                if j == 0:
                    p.font.bold = True
                    p.alignment = PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1a, 0x23, 0x40) if i % 2 == 0 else RGBColor(0x22, 0x2b, 0x50)

    # Observations
    lines = [
        ("Key observations:", True, ACCENT),
        ("- Room experiment shows highest bg preservation (0.924 CLIP) -- simpler prompts help", False, WHITE),
        ("- Horse experiment has strongest fg_clip_score (0.120) -- clearest object transformation", False, WHITE),
        ("- Negative fg_clip for room suggests delta embedding poorly captures table->aquarium change", False, WHITE),
    ]
    add_multiline(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.5), lines, font_size=15)


def make_conclusions_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.8),
             "Summary & Next Steps", font_size=36, bold=True, color=ACCENT)

    lines = [
        ("Results:", True, ACCENT),
        ("- Full enumeration of 2^20 binary diffusion paths for 3 different editing tasks", False, WHITE),
        ("- DINOv2 embeddings reveal structured clusters in the path space", False, WHITE),
        ("- UMAP + HDBSCAN discovers 50-65 distinct visual clusters per experiment", False, WHITE),
        ("- Metrics (bg_clip, bg_ssim, fg_clip) allow automated ranking of edit quality", False, WHITE),
        ("", False, None),
        ("Observations:", True, ACCENT),
        ("- Path space has clear structure: nearby binary codes -> similar images", False, WHITE),
        ("- Combined metrics identify best edits: strong foreground change + preserved background", False, WHITE),
        ("- Different editing tasks produce different cluster topologies", False, WHITE),
        ("", False, None),
        ("Next steps:", True, ORANGE),
        ("- Analyze bit importance: which diffusion steps matter most for editing?", False, WHITE),
        ("- Explore interpolation strategies beyond binary (ternary, continuous)", False, WHITE),
        ("- Scale to more complex editing scenarios", False, WHITE),
    ]
    add_multiline(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), lines, font_size=18)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    A = f"{PROJ}/analysis"

    # 1. Title
    make_title_slide(prs)

    # 2. Method
    make_method_slide(prs)

    # --- CATDOG ---
    make_experiment_slide(
        prs,
        title="Experiment 1: Cat -> Dog",
        source_prompt="tabby kitten walking confidently across a stone pavement.",
        target_prompt="tabby dog walking confidently across a stone pavement.",
        seg_prompt="tabby kitten",
        source_img=f"{PROJ}/generation/istanbul-cats-history.jpg",
        path_source=f"{NFS3}/generated_samples_40step/path_00000_b0.jpg",
        path_target=f"{NFS3}/generated_samples_40step/path_00001_b1048575.jpg",
        seg_vis=f"{PROJ}/metrics/background_mask_vis.jpg",
        metrics_text=[
            ("Metrics (1,048,576 images):", True, ACCENT),
            ("bg_clip_similarity: 0.848 +/- 0.059    bg_ssim: 0.545 +/- 0.024", False, WHITE),
            ("fg_clip_score: 0.045 +/- 0.067          HDBSCAN clusters: 52", False, WHITE),
        ],
        num_images="1,048,576",
    )
    make_clusters_slide(prs, "Cat -> Dog",
                        f"{A}/catdog_analysis/map_clusters.png",
                        f"{A}/catdog_analysis/map_seg_metrics.png")
    make_grid_slide(prs, "Cat -> Dog", f"{A}/catdog_analysis/grid_top_combined.png")

    # --- HORSE ---
    make_experiment_slide(
        prs,
        title="Experiment 2: Horse -> Robot Horse",
        source_prompt="a horse on the grass",
        target_prompt="a robot horse on the grass",
        seg_prompt="horse",
        source_img=f"{PROJ}/generation/horse.jpg",
        path_source=f"{NFS3}/generated_horse_300k/path_00000_b0.jpg",
        path_target=f"{NFS3}/generated_horse_300k/path_00001_b1048575.jpg",
        seg_vis=f"{PROJ}/metrics/background_mask_horse_vis.jpg",
        metrics_text=[
            ("Metrics (1,048,576 images):", True, ACCENT),
            ("bg_clip_similarity: 0.862 +/- 0.047    bg_ssim: 0.623 +/- 0.042", False, WHITE),
            ("fg_clip_score: 0.120 +/- 0.051          HDBSCAN clusters: 53", False, WHITE),
        ],
        num_images="1,048,576",
    )
    make_clusters_slide(prs, "Horse -> Robot",
                        f"{A}/horse_analysis/map_clusters.png",
                        f"{A}/horse_analysis/map_seg_metrics.png")
    make_grid_slide(prs, "Horse -> Robot", f"{A}/horse_analysis/grid_top_combined.png")

    # --- ROOM ---
    make_experiment_slide(
        prs,
        title="Experiment 3: Coffee Table -> Aquarium",
        source_prompt="A coffee table and a sofa in a modern living room.",
        target_prompt="An aquarium and a sofa in a modern living room.",
        seg_prompt="coffee table",
        source_img=f"{PROJ}/generation/room.png",
        path_source=f"{A}/room_analysis/path_source_b0.jpg",
        path_target=f"{A}/room_analysis/path_target_b1048575.jpg",
        seg_vis=f"{PROJ}/metrics/background_mask_room_vis.jpg",
        metrics_text=[
            ("Metrics (1,048,576 images):", True, ACCENT),
            ("bg_clip_similarity: 0.924 +/- 0.027    bg_ssim: 0.572 +/- 0.073", False, WHITE),
            ("fg_clip_score: -0.156 +/- 0.019         HDBSCAN clusters: 65", False, WHITE),
        ],
        num_images="1,048,576",
    )
    make_clusters_slide(prs, "Table -> Aquarium",
                        f"{A}/room_analysis/map_clusters.png",
                        f"{A}/room_analysis/map_seg_metrics.png")
    make_grid_slide(prs, "Table -> Aquarium", f"{A}/room_analysis/grid_top_combined_ssim.png")

    # --- COMPARISON ---
    make_comparison_slide(prs)

    # --- CONCLUSIONS ---
    make_conclusions_slide(prs)

    out = f"{PROJ}/analysis/presentation_results.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
