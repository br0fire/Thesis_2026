"""
Fill the Skoltech ACS/DS MSc pre-defense template with thesis content.

Takes the stock template, keeps all styling (fonts, colours, Skoltech logo,
slide layouts), and replaces the boilerplate body text on each of the 15
slides with research content. Output: presentation_pptx.pptx.
"""
from copy import deepcopy
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import os, shutil

SRC = "/home/jovyan/shares/SR006.nfs2/svgrozny/project/clear_project/docs/presentation/TSR, Pre-Defense, Defense PPT template - ACS & DS MSc.pptx"
OUT = "/home/jovyan/shares/SR006.nfs2/svgrozny/project/clear_project/docs/presentation/presentation_pptx.pptx"

THESIS_TITLE = "Exploring Diffusion Trajectories"
STUDENT      = "Sergey Grozny"
ADVISOR      = "Aleksandr Katrutsa, PhD"
SHORT_TITLE  = "Grozny. Exploring Diffusion Trajectories"
PROGRAM      = "Data Science, MSc — Pre-defense"

FIG_DIR = "/home/jovyan/shares/SR006.nfs3/svgrozny/clean_v1"
THESIS_FIG = "/home/jovyan/shares/SR006.nfs2/svgrozny/project/clear_project/docs/Skoltech_MSc_Thesis_Template_DS/images/phase2"


def clear_frame(tf):
    """Remove all paragraphs from a text frame; return the saved first-run font settings."""
    # capture font size & color from first run if any
    size, color_rgb, font_name = None, None, None
    for para in tf.paragraphs:
        for r in para.runs:
            if r.font.size is not None:
                size = r.font.size
            if r.font.name is not None:
                font_name = r.font.name
            try:
                if r.font.color and r.font.color.rgb is not None:
                    color_rgb = r.font.color.rgb
            except Exception:
                pass
            break
        if size is not None:
            break
    # remove all <a:p> children
    tfxml = tf._txBody
    for p in tfxml.findall(qn("a:p")):
        tfxml.remove(p)
    return size, color_rgb, font_name


def write_lines(tf, lines, default_size_pt=20, bold_first=False, bullet=True):
    """Write a sequence of lines (strings or (text, bold) tuples) into an empty text frame."""
    size, color_rgb, font_name = clear_frame(tf)
    if size is None:
        size = Pt(default_size_pt)

    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, is_bold = line
        else:
            text, is_bold = line, (bold_first and i == 0)

        p = tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = text
        r.font.size = size
        r.font.bold = is_bold
        if font_name:
            r.font.name = font_name
        if color_rgb is not None:
            r.font.color.rgb = color_rgb


def set_placeholder_text(slide, idx, lines, default_size_pt=20, bullet=True, bold_first=False):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            write_lines(ph.text_frame, lines, default_size_pt=default_size_pt, bullet=bullet, bold_first=bold_first)
            return ph
    return None


def set_footer(slide, text=SHORT_TITLE):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 11:
            write_lines(ph.text_frame, [text])
            return


def add_image(slide, path, left_in, top_in, width_in=None, height_in=None):
    from pptx.util import Inches
    kw = {}
    if width_in is not None:
        kw["width"] = Inches(width_in)
    if height_in is not None:
        kw["height"] = Inches(height_in)
    return slide.shapes.add_picture(path, Inches(left_in), Inches(top_in), **kw)


def main():
    shutil.copy(SRC, OUT)
    prs = Presentation(OUT)
    slides = list(prs.slides)

    # ─── Slide 1 — Cover ──────────────────────────────────────────────────
    s = slides[0]
    # idx=0 title
    set_placeholder_text(s, 0, [THESIS_TITLE], default_size_pt=40, bold_first=True)
    # idx=1 program
    set_placeholder_text(s, 1, [PROGRAM], default_size_pt=26)
    # idx=3 student/advisor
    set_placeholder_text(s, 3, [
        f"Student: {STUDENT}",
        f"Research Advisor: {ADVISOR}",
        "",
    ], default_size_pt=26)

    # ─── Slide 2 — General problem ────────────────────────────────────────
    s = slides[1]
    set_placeholder_text(s, 0, ["Problem: training-free image editing"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Text-to-image diffusion models produce photorealistic outputs from a prompt, but editing an existing generation requires preserving unrelated content (background, layout, lighting) while changing only what the new prompt specifies.",
        "",
        "Existing training-free approaches each pay a cost:",
        "   •  SDEdit — single noise-level knob, coarse.",
        "   •  Prompt-to-Prompt — attention-map surgery; needs model internals.",
        "   •  DreamBooth / LoRA — per-edit fine-tuning; expensive and inflexible.",
        "",
        "Observation: denoising runs T timesteps; the text prompt is usually fixed across all of them. Early steps commit layout, late steps resolve details (Choi et al., 2022). What if we switch the prompt per timestep?",
    ], default_size_pt=20)
    set_footer(s)

    # ─── Slide 3 — Aim ────────────────────────────────────────────────────
    s = slides[2]
    set_placeholder_text(s, 0, ["Aim"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Build a training-free image-editing framework that treats an edit as combinatorial search over binary per-step prompt-conditioning schedules — without fine-tuning the backbone, without access to model internals, and at a practical sampling budget.",
        "",
        "Potential impact:",
        "   •  Editing over black-box diffusion APIs (no gradients, no weights exposed).",
        "   •  Backbone-agnostic recipe that ports from Stable Diffusion v1.5 to FLUX.2 without code changes.",
        "   •  Two orders of magnitude cheaper than exhaustive enumeration at matched quality.",
    ], default_size_pt=20)
    set_footer(s)

    # ─── Slide 4 — Objectives ─────────────────────────────────────────────
    s = slides[3]
    set_placeholder_text(s, 0, ["Objectives"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "1.  Formalise editing as a search over binary conditioning masks M ∈ {0,1}ᵀ (Discrete Context Switching).",
        "2.  Map the structure of the 2ᵀ trajectory space by exhaustive enumeration on SD v1.5 (2²⁰ masks per task).",
        "3.  Design a segmentation-aware reward that separately scores background preservation and foreground editability.",
        "4.  Replace enumeration with a REINFORCE policy over 14 Bernoulli logits on FLUX.2; validate on 15 editing tasks.",
        "5.  Compare against a matched-budget random baseline and the true 2¹⁴ exhaustive ceiling.",
        "6.  Audit the algorithm landscape: evolutionary, hill-climb, CEM, Thompson, with and without priors.",
    ], default_size_pt=20)
    set_footer(s)

    # ─── Slide 5 — Theory / Algorithms ────────────────────────────────────
    s = slides[4]
    set_placeholder_text(s, 0, ["Theory and algorithms"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Discrete Context Switching (DCS).  A binary mask M = [mᴛ, …, m₁] ∈ {0,1}ᵀ selects source or target conditioning per denoising step:",
        "       c_t = m_t · ψ(P_tgt) + (1 − m_t) · ψ(P_src).",
        "Hypothesis space |H| = 2ᵀ. Default baseline M = 1 (pure-target prompt every step).",
        "",
        "Segmentation-aware reward.  R(M) = f_bg(M)^α · ( relu(f_fg(M)) + ε )^(1−α), α = 0.5 – 0.7.",
        "       f_bg — SSIM restricted to bg pixels from SAM 3.1.     f_fg — cosine of SigLIP-2 crop embedding with normalised Δ_text = ψ(P_tgt) − ψ(P_src).",
        "",
        "Policy.  π_θ(M) = ∏ᵢ Bern(mᵢ | σ(θᵢ)), 14 independent logits.",
        "REINFORCE update with EMA baseline, per-batch advantage normalisation, and entropy regularisation (β = 0.05).",
        "Early stop: Kendall-τ plateau on mean reward  OR  entropy H < 0.5. Floor 200 episodes, cap 300 (80 in fixed-budget sweeps).",
    ], default_size_pt=16)
    set_footer(s)

    # ─── Slide 6 — Methodology / Experimental setup ───────────────────────
    s = slides[5]
    set_placeholder_text(s, 0, ["Methodology and experimental setup"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Backbones.  (Phase 1) Stable Diffusion v1.5 with Null-Text Inversion, DDIM 40 steps. (Phase 2) FLUX.2-klein-base-9B, rectified flow, 28 steps, inversion-free (source prompt + fixed seed gives the source image).",
        "",
        "Reward model.  SigLIP-2 so400m-patch14-384 (smoother cosine landscape than CLIP ViT-B/32 on cropped crops).",
        "Segmentation.  SAM 3.1 → fallback Grounding-DINO + SAM 2 → fallback CLIPSeg.",
        "",
        "Tasks.  Phase 1: 2 probes (cat→dog, horse→robot horse), 2²⁰ masks each. Phase 2: 15 background-rich edits at 14-bit / 640 generations, plus 8 scenes with full 2¹⁴ enumeration for ceiling analysis.",
        "",
        "Hardware.  8×A100 80 GB for Phase-1 sweep; single A100 per Phase-2 edit. Canonical source.pt + bg_mask.npy reused across methods (avoids CUDA non-determinism).",
        "Storage.  Large tensors on NFS-3; code on NFS-2; memmap for concurrent writes.",
    ], default_size_pt=16)
    set_footer(s)

    # ─── Slide 7 — Results ────────────────────────────────────────────────
    s = slides[6]
    set_placeholder_text(s, 0, ["Results"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Phase 1 (SD v1.5, exhaustive 2²⁰).  HDBSCAN on DINOv2+UMAP gives 52 clusters for cat→dog, 13 for horse→robot horse. All-ones reward 0.220 / 0.167 vs best-mask 0.885 / 0.826; over 95% of random masks beat M = 1 on both tasks.",
        "",
        "Phase 2 (FLUX.2, REINFORCE 80 ep × 8 = 640 gen).  Mean reward climbs 0.601 → 0.749 (+23%, +0.148) across all 15 background-rich edits. Wins vs all-ones on 15 / 15. 80% of final best-of-run reached inside the first 10 episodes.",
        "",
        "Ceiling analysis (8 canonical edits with full 2¹⁴ enumeration):",
        "   •  REINFORCE @ 640 reaches  99.6 – 100 %  of the 2¹⁴ exhaustive maximum.",
        "   •  Random-640 reaches       95 – 100 %; the algorithmic edge is +0.6 – 2.3 % best-of-N.",
        "   •  Framework win (+0.148 vs M = 1) comes from path-mixing itself, not the specific search algorithm.",
        "",
        "Algorithm audit (lookup sweep, 550+ variants).  Without prior: ev_pop32 = 0.9988, sim_anneal = 0.9970, REINFORCE = 0.9955, random = 0.9918. With leave-one-out bgrich prior: hill_prior reaches 96.4% at budget = 1.",
    ], default_size_pt=14)
    set_footer(s)

    # ─── Slide 8 — Discussion ─────────────────────────────────────────────
    s = slides[7]
    set_placeholder_text(s, 0, ["Discussion"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "The source-early / target-late pattern is backbone-independent. It shows up in SD v1.5 exhaustive (2²⁰ masks, P2P-style schedule) and in FLUX.2 REINFORCE logits (different backbone, different reward model, different optimiser) — evidence that it is a property of diffusion trajectories, not of one architecture.",
        "",
        "The reward landscape on 14-bit bg-rich scenes is wide-flat in the top-1 %.  Random-640 catches up because 640 uniform samples cover 3.9 % of the 2¹⁴ space and hit the plateau by birthday paradox. REINFORCE's algorithmic edge is real but small (1 – 3 %); the main win is vs the all-ones baseline, not vs random.",
        "",
        "The bgrich prior does not transfer.  On 8 diverse edits (object / landscape / portrait / style) at matched budget 120, prior-REINFORCE loses to random-120 on 8/8 and to no-prior REINFORCE on 5/8 (mean 0.594 vs 0.645 vs 0.709). The prior is distribution-specific to LLM-template 'object on rich background' scenes, not a universal recommendation.",
        "",
        "Personal contribution.  DCS formulation, SD-v1.5 exhaustive pipeline, segmentation-aware reward, REINFORCE-on-FLUX code, 15-edit benchmark, ceiling analysis, diverse-prior generalisation test.",
    ], default_size_pt=15)
    set_footer(s)

    # ─── Slide 9 — Scientific novelty ─────────────────────────────────────
    s = slides[8]
    set_placeholder_text(s, 0, ["Scientific novelty"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "1.  Discrete Context Switching — a binary-mask reformulation of diffusion editing.  No weight updates, no inversion, no attention-internal surgery. Strictly more general than P2P's single threshold (accepts non-monotonic schedules).",
        "",
        "2.  Backbone-portable editing recipe.  Same framework runs on SD v1.5 (DDIM + Null-Text Inversion) and FLUX.2 (rectified flow, inversion-free). The source-early / target-late pattern replicates across both.",
        "",
        "3.  First 2²⁰ exhaustive diffusion-editing map.  UMAP + HDBSCAN clustering of ~10⁶ generations per task reveals 52 / 13 geometrically distinct clusters and the M = 1 trajectory sits outside the top cluster on both probes.",
        "",
        "4.  REINFORCE on a 14-bit Bernoulli-logit policy with plateau early-stop reaches 99.6 – 100 % of the 2¹⁴ exhaustive ceiling at 640 generations — a 25× budget reduction with no quality loss against matched-budget random.",
        "",
        "5.  Honest algorithm audit: shows that on a wide-flat landscape the algorithm matters less than the framework. Warns that priors derived from LLM-template prompts do not transfer.",
    ], default_size_pt=14)
    set_footer(s)

    # ─── Slide 10 — Innovation ────────────────────────────────────────────
    s = slides[9]
    set_placeholder_text(s, 0, ["Innovation"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Industrial application.  DCS is a practical editing backend for black-box diffusion APIs — the model weights do not need to be exposed, only the ability to pass a per-step text embedding. Works unchanged for any API that offers prompt-per-step control (already true for FLUX.2 and straightforward to add to SD-class servers).",
        "",
        "Amortised deployment.  Once a small bank of 30+ completed edits is collected, a prompt-conditioned regressor (SigLIP embeddings → 14 Bernoulli logits) predicts a near-optimal greedy mask in one forward pass. Inference collapses to one diffusion call at quality within 4% of the exhaustive ceiling on matching prompt distributions.",
        "",
        "Deployment cost.  640 generations per edit ≈ 2 wall-clock hours on a single A100 in Phase-2 settings; amortised mode ≈ 1 diffusion call, i.e. real-time.",
    ], default_size_pt=16)
    set_footer(s)

    # ─── Slide 11 — Conclusions ───────────────────────────────────────────
    s = slides[10]
    set_placeholder_text(s, 0, ["Conclusions"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "The binary-mask trajectory space contains edits that outperform the all-target baseline.  Confirmed by exhaustive 2²⁰ search on SD v1.5 (> 95% of random masks beat M = 1) and by all-ones comparison on 15 FLUX.2 edits (+0.148 reward / +23%).",
        "",
        "Top-reward trajectories activate source early, target late.  Pattern holds across backbone (SD v1.5 vs FLUX.2), reward model (CLIP ViT-B/32 vs SigLIP-2 so400m), and optimisation regime (exhaustive argmax vs REINFORCE logits).",
        "",
        "The segmentation-aware bg / fg reward is a viable scalar signal for gradient-free policy optimisation on a black-box generator.  It is differentiable through the policy, not through the diffusion model.",
        "",
        "REINFORCE over 14 Bernoulli logits with entropy regularisation and plateau early-stop makes DCS practical — 640 generations reach 99.6 – 100 % of the 2¹⁴ exhaustive ceiling without losing quality to matched-budget random. The framework is significant for production editing pipelines where fine-tuning budgets or backbone access are unavailable.",
    ], default_size_pt=15)
    set_footer(s)

    # ─── Slide 12 — Current Status ────────────────────────────────────────
    s = slides[11]
    set_placeholder_text(s, 0, ["Current status"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Thesis writing.  Draft is complete through Chapter 8 (Conclusion). Abstract, intro, motivation, problem statement, methods, numerical experiments chapters are written; defense statements finalised.",
        "",
        "Experiments done.",
        "   •  Phase-1 exhaustive 2²⁰ on SD v1.5: cat→dog + horse→robot horse, full UMAP+HDBSCAN analysis.",
        "   •  Phase-2 REINFORCE on FLUX.2: 15 background-rich edits at 80 episodes, 8-config hyperparameter sweep.",
        "   •  Clean-v1 canonical-source pipeline: 8 canonical sources, 2¹⁴ exhaustive ceilings, ceiling-ratio analysis.",
        "   •  Search-algorithm audit (550+ variants on CPU lookup).",
        "   •  Diverse-prior generalisation test (8 edits, 3 methods each).",
        "",
        "Open.  Cross-distribution validation at matched 640-episode budget; amortised MLP/Ridge policy from 30+ edits; perceptual / DINOv2-based background reward.",
    ], default_size_pt=15)
    set_footer(s)

    # ─── Slide 13 — Outcomes ──────────────────────────────────────────────
    s = slides[12]
    set_placeholder_text(s, 0, ["Outcomes"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Code.  Open-source release planned alongside the thesis: generation/ pipeline, analysis/ sweep and aggregation scripts, 16-prompt benchmark fixtures.",
        "",
        "Data.  2²⁰ SD-v1.5 exhaustive outputs + metrics for both probes; 2¹⁴ FLUX.2 exhaustive outputs + metrics for 8 canonical bg-rich edits.",
        "",
        "Artifacts for the defense.",
        "   •  8-panel training-curves figure (mean reward, best-so-far, bg SSIM, fg CLIP, entropy).",
        "   •  Per-experiment 4-metric grids with exhaustive ceiling and all-ones reference lines.",
        "   •  Reward-landscape histograms showing where REINFORCE shifts the mass vs random-640.",
        "   •  Visual comparison grid: source / all-ones / best / amortised greedy.",
        "",
        "Planned paper: short workshop note (\"Source-early / target-late is a diffusion-editing invariant\") submitted to an NeurIPS / CVPR workshop cycle after thesis submission.",
    ], default_size_pt=15)
    set_footer(s)

    # ─── Slide 14 — Outlook ───────────────────────────────────────────────
    s = slides[13]
    set_placeholder_text(s, 0, ["Outlook"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Cross-distribution prior validation.  The bgrich prior fails to transfer — next step is a 640-episode run on the 8 diverse scenes and a mixed-distribution prior pooled from bgrich + diverse exhaustives.",
        "",
        "Peakier rewards.  Wide-flat top-1% landscape caps any algorithm's edge over random. Candidates: perceptual similarity (LPIPS) or DINOv2-feature cosine for background; VLM-scored crop for foreground.",
        "",
        "Amortised policy.  15 completed REINFORCE runs are not enough (MLP loses to population mean). Overnight Phase-3 expanding to 30 runs; Ridge + SigLIP inputs are the lowest-risk baseline.",
        "",
        "n_bits = 28 regime.  28-bit (no step repetition) opens trajectories unreachable at 14-bit (+0.015 reward). Enumeration is infeasible (2²⁸ = 268M) — only sampling methods will work; REINFORCE / CEM / ev_pop32 under this regime is the natural next test.",
        "",
        "Larger-scale study.  8-experiment benchmark is statistically tight; scaling to 30 – 50 would tighten confidence intervals but costs proportional FLUX time.",
    ], default_size_pt=14)
    set_footer(s)

    # ─── Slide 15 — Acknowledgements ──────────────────────────────────────
    s = slides[14]
    set_placeholder_text(s, 0, ["Acknowledgements"], default_size_pt=32, bold_first=True)
    set_placeholder_text(s, 1, [
        "Aleksandr Katrutsa (research advisor) — project direction, technical review, weekly guidance.",
        "",
        "Skoltech AI-R group — shared compute, infrastructure support.",
        "",
        "Compute.  8× NVIDIA A100 80 GB cluster, Skoltech MLSpace.",
        "",
        "Open-source models and code: FLUX.2-klein (Black Forest Labs), Stable Diffusion v1.5 (CompVis), SAM 3 (Meta), SigLIP-2 (Google), Grounding DINO, CLIPSeg, DINOv2.",
    ], default_size_pt=18)
    set_footer(s)

    prs.save(OUT)
    print(f"Wrote: {OUT}  ({os.path.getsize(OUT)/1024:.0f} KB)")


if __name__ == "__main__":
    main()
